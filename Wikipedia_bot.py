import speech_recognition as sr
import pyttsx3
import torch
from transformers import Wav2Vec2ForCTC, Wav2Vec2Tokenizer
import soundfile as sf
import string
import json
import os
from tkinter import Tk
from tkinter.filedialog import askopenfilename
import PyPDF2
import docx
from pptx import Presentation
import datetime
import wikipedia
import wikipediaapi

# Initialize Wikipedia API with a proper User-Agent
wiki_wiki = wikipediaapi.Wikipedia(
    language='en',
    user_agent='WikipediaBot/1.0 (https://example.com; contact@example.com)'
)

# Configure settings for the `wikipedia` library
wikipedia.set_lang("en")
wikipedia.set_rate_limiting(True)

# Load pre-trained Wav2Vec2 model and tokenizer for offline STT
tokenizer = Wav2Vec2Tokenizer.from_pretrained("facebook/wav2vec2-base-960h")
model = Wav2Vec2ForCTC.from_pretrained("facebook/wav2vec2-base-960h")

# Initialize text-to-speech engine (pyttsx3)
engine = pyttsx3.init()

# Path to the knowledge base JSON file
knowledge_base_file = r"C:\Users\noher\Downloads\chat_formatted1.json"

# Load the knowledge base from the JSON file
def load_knowledge_base():
    if os.path.exists(knowledge_base_file):
        with open(knowledge_base_file, 'r') as f:
            return json.load(f)
    else:
        return {}

# Save the updated knowledge base to the JSON file
def save_knowledge_base():
    with open(knowledge_base_file, 'w') as f:
        json.dump(knowledge_base, f, indent=4)

# Initialize knowledge base from file
knowledge_base = load_knowledge_base()

# ------------------ Text Preprocessing Functions ------------------

# Tokenizer - Split the input into words (tokens)
def tokenize(text):
    return text.lower().split()

# Normalization - Remove punctuation
def normalize(tokens):
    table = str.maketrans('', '', string.punctuation)
    return [word.translate(table) for word in tokens]

# Stop Words Removal - Ignore common words like "the", "is", etc.
stop_words = {'the', 'is', 'at', 'on', 'and', 'a', 'in', 'of', 'to'}
def remove_stop_words(tokens):
    return [word for word in tokens if word not in stop_words]

# ------------------ Date and Time Functionality ------------------

# Function to get the current date and time
def get_current_datetime():
    now = datetime.datetime.now()
    return now.strftime("%Y-%m-%d %H:%M:%S")

# Function to handle date and time queries
def handle_date_time_query(user_input):
    if "current time" in user_input or "time now" in user_input:
        current_time = datetime.datetime.now().strftime("%H:%M:%S")
        return f"The current time is {current_time}."
    elif "current date" in user_input or "today's date" in user_input:
        current_date = datetime.datetime.now().strftime("%Y-%m-%d")
        return f"Today's date is {current_date}."
    else:
        return None

# ------------------ Interaction Logging ------------------

# Path to interaction log file
log_file = "interaction_log.txt"

# Log interaction with timestamps
def log_interaction(user_input, bot_response):
    timestamp = get_current_datetime()
    with open(log_file, 'a') as f:
        f.write(f"{timestamp} - User: {user_input}\n")
        f.write(f"{timestamp} - Bot: {bot_response}\n")

# ------------------ Voice Input (STT) and Output (TTS) Functions ------------------

# Speech-to-Text (using Hugging Face's Wav2Vec2 for offline functionality)
def read_wav_file(wav_file):
    speech, sample_rate = sf.read(wav_file)
    return speech

def wav_to_text(wav_file):
    speech = read_wav_file(wav_file)
    input_values = tokenizer(speech, return_tensors="pt", padding="longest").input_values
    with torch.no_grad():
        logits = model(input_values).logits

    predicted_ids = torch.argmax(logits, dim=-1)
    transcription = tokenizer.decode(predicted_ids[0])
    return transcription.lower()

# Speech-to-text using Google's API (fallback if Wav2Vec2 isn't used)
def speech_to_text():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening...")
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source)
   
    try:
        text = recognizer.recognize_google(audio)
        print(f"User said: {text}")
        return text.lower()
    except sr.UnknownValueError:
        print("Sorry, I did not catch that.")
        return None
    except sr.RequestError as e:
        print(f"Could not request results; {e}")
        return None

# Text-to-speech using pyttsx3
def text_to_speech(text):
    engine.say(text)
    engine.runAndWait()

# ------------------ File Upload Functions ------------------

# Open file dialog to allow file upload
def choose_file():
    root = Tk()
    root.withdraw()  # Hide the main window
    root.update()  # Ensure dialog box opens properly
    
    file_path = askopenfilename(title="Select a file", filetypes=[("All files", "*.*"), 
                                                                ("PDF files", "*.pdf"), 
                                                                ("DOCX files", "*.docx"), 
                                                                ("PPTX files", "*.pptx")])
    
    if file_path:
        print(f"File chosen: {file_path}")
        return file_path
    else:
        print("No file chosen.")
        return None

# Read the contents of the file
def read_file(file_path):
    file_ext = os.path.splitext(file_path)[-1].lower()
    
    if file_ext == '.txt':
        with open(file_path, 'r') as f:
            content = f.read()
    elif file_ext == '.pdf':
        content = read_pdf(file_path)
    elif file_ext == '.docx':
        content = read_docx(file_path)
    elif file_ext == '.pptx':
        content = read_pptx(file_path)
    else:
        content = f"Unsupported file type: {file_ext}. Supported formats are .txt, .pdf, .docx, .pptx."
    
    return content

# Read PDF content
def read_pdf(file_path):
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        content = ""
        for page_num in range(len(reader.pages)):
            content += reader.pages[page_num].extract_text()
    return content

# Read DOCX content
def read_docx(file_path):
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content

# Read PPTX content
def read_pptx(file_path):
    presentation = Presentation(file_path)
    content = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content

# ------------------ Wikipedia Functionality ------------------

# Enhanced function to fetch content from Wikipedia
def search_wikipedia(query):
    # Try to fetch using wikipediaapi
    page = wiki_wiki.page(query)
    if page.exists():
        # Return summary from wikipediaapi
        summary = page.summary[:500] + ("..." if len(page.summary) > 500 else "")
        return summary
    
    # Fallback to the `wikipedia` library for full-text search
    try:
        print("Using fallback Wikipedia search...")
        search_results = wikipedia.search(query)
        if search_results:
            # Use the first result and fetch summary
            page_title = search_results[0]
            page_summary = wikipedia.summary(page_title, sentences=2)
            return page_summary
        else:
            return "I couldn't find anything on Wikipedia for that topic. Could you provide more details?"
    except wikipedia.exceptions.DisambiguationError as e:
        # Handle ambiguous queries
        return f"Your query is ambiguous. Did you mean: {', '.join(e.options[:5])}?"
    except wikipedia.exceptions.PageError:
        return "I couldn't find anything on Wikipedia for that topic."
    except Exception as e:
        # General exception handling
        print(f"Error accessing Wikipedia: {e}")
        return "An error occurred while trying to fetch information from Wikipedia."

# ------------------ Conversational AI and Learning Functions ------------------

# Response generation from knowledge base or Wikipedia
def generate_response_with_wikipedia(user_input):
    tokens = tokenize(user_input)
    normalized_tokens = normalize(tokens)
    cleaned_input = ' '.join(normalized_tokens)

    if "upload a file" in user_input or "i want to upload a file" in user_input:
        file_path = choose_file()  # Open file dialog
        if file_path:
            file_content = read_file(file_path)  # Read the file content
            knowledge_base[file_path] = file_content  # Store file content in knowledge base
            save_knowledge_base()
            text_to_speech("File uploaded and stored in the knowledge base.")
            return "File uploaded successfully!"
        else:
            return "No file uploaded."

    date_time_response = handle_date_time_query(user_input)
    if date_time_response:
        return date_time_response

    # Check knowledge base
    if cleaned_input in knowledge_base:
        return knowledge_base[cleaned_input]
    else:
        # Enhanced Wikipedia search
        print("Searching Wikipedia for an answer...")
        wikipedia_response = search_wikipedia(user_input)
        if "I couldn't find anything on Wikipedia" not in wikipedia_response:
            # Add to knowledge base
            knowledge_base[cleaned_input] = wikipedia_response
            save_knowledge_base()
        return wikipedia_response

# ------------------ Main Conversation Engine ------------------

def conversation_engine_with_voice():
    print("Bot: Hello! Let's chat. Speak to me. Say 'bye' to exit.")
    text_to_speech("Hello! Let's chat. Speak to me. Say 'bye' to exit.")
   
    while True:
        # Capture speech input
        user_input = speech_to_text()
        if not user_input:
            print("Could not understand. Please try again.")
            continue
       
        print(f"You said: {user_input}")
       
        if user_input == 'bye':
            print("Bot: Goodbye!")
            text_to_speech("Goodbye!")
            break
       
        # Generate response with enhanced Wikipedia fallback
        bot_response = generate_response_with_wikipedia(user_input)
        print(f"Bot: {bot_response}")
       
        # Convert response to speech
        text_to_speech(bot_response)
       
        # Log interaction
        log_interaction(user_input, bot_response)

# Start the conversation with voice interaction
if __name__ == "__main__":
    conversation_engine_with_voice()